from __future__ import annotations

import contextlib
import json
import markdown as _md_lib
import re
import threading
import time
from datetime import datetime
from pathlib import Path
from urllib.parse import urlparse

import arxiv
import httpx
from bs4 import BeautifulSoup
from ddgs import DDGS
from langchain_community.tools import WikipediaQueryRun
from langchain_community.utilities import WikipediaAPIWrapper
from langchain_core.messages import HumanMessage, SystemMessage
from langchain_core.tools import tool
from youtube_transcript_api import NoTranscriptFound, TranscriptsDisabled, YouTubeTranscriptApi

from ..settings import get_user_tool_settings
from ..core import activity_log as actlog
from ..core.activity_log import EVT_TOOL_RETRY, EVT_TOOL_TIMEOUT
from .llm_factory import build_summary_llms, invoke_with_fallbacks


_tls = threading.local()
_wiki = WikipediaQueryRun(api_wrapper=WikipediaAPIWrapper(top_k_results=1, doc_content_chars_max=1500))

# ── visit_url constants ───────────────────────────────────────────────────────
_URL_FETCH_TIMEOUT = 15
_URL_MAX_CHARS = 8_000
_URL_MAX_REDIRECTS = 5
_URL_RETRY_ATTEMPTS = 3
_URL_RETRY_BACKOFF_BASE_SECS = 0.4
_URL_MAX_CONCURRENCY_PER_DOMAIN = 2
_URL_RETRYABLE_STATUS_CODES = {429, 500, 502, 503, 504}
_URL_HEADERS = {
    "User-Agent": (
        "Mozilla/5.0 (Windows NT 10.0; Win64; x64) "
        "AppleWebKit/537.36 (KHTML, like Gecko) "
        "Chrome/122.0.0.0 Safari/537.36"
    ),
    "Accept": "text/html,application/xhtml+xml,application/xml;q=0.9,*/*;q=0.8",
    "Accept-Language": "en-US,en;q=0.5",
}


class _DomainConcurrencyLimiter:
    """Bound concurrent outbound requests per domain."""

    def __init__(self, per_domain_limit: int) -> None:
        self._per_domain_limit = max(1, int(per_domain_limit))
        self._lock = threading.Lock()
        self._semaphores: dict[str, threading.BoundedSemaphore] = {}

    def _get_sem(self, domain: str) -> threading.BoundedSemaphore:
        with self._lock:
            sem = self._semaphores.get(domain)
            if sem is None:
                sem = threading.BoundedSemaphore(self._per_domain_limit)
                self._semaphores[domain] = sem
            return sem

    @contextlib.contextmanager
    def slot(self, domain: str):
        sem = self._get_sem(domain)
        sem.acquire()
        try:
            yield
        finally:
            sem.release()


_visit_url_domain_limiter = _DomainConcurrencyLimiter(per_domain_limit=_URL_MAX_CONCURRENCY_PER_DOMAIN)


def _domain_key_from_url(url: str) -> str:
    try:
        host = urlparse(url).netloc.lower().strip()
    except Exception:
        host = ""
    if host.startswith("www."):
        host = host[4:]
    return host or "unknown-domain"


def _fetch_url_with_retries(url: str) -> httpx.Response:
    """Robust URL fetch with per-domain concurrency, redirect cap, and retry/backoff."""
    domain = _domain_key_from_url(url)
    timeout = httpx.Timeout(_URL_FETCH_TIMEOUT, connect=_URL_FETCH_TIMEOUT)
    last_exc: Exception | None = None

    for attempt in range(1, _URL_RETRY_ATTEMPTS + 1):
        try:
            with _visit_url_domain_limiter.slot(domain):
                with httpx.Client(
                    headers=_URL_HEADERS,
                    timeout=timeout,
                    follow_redirects=True,
                    max_redirects=_URL_MAX_REDIRECTS,
                ) as client:
                    response = client.get(url)

            if response.status_code in _URL_RETRYABLE_STATUS_CODES:
                if attempt < _URL_RETRY_ATTEMPTS:
                    delay = _URL_RETRY_BACKOFF_BASE_SECS * (2 ** (attempt - 1))
                    actlog.log(EVT_TOOL_RETRY, {
                        "url": url, "attempt": attempt,
                        "reason": f"http_{response.status_code}", "backoff_s": delay
                    })
                    time.sleep(delay)
                    continue
                response.raise_for_status()

            response.raise_for_status()
            return response

        except (httpx.TimeoutException, httpx.RequestError) as exc:
            last_exc = exc
            if isinstance(exc, httpx.TimeoutException):
                actlog.log(EVT_TOOL_TIMEOUT, {
                    "url": url, "attempt": attempt, "error": str(exc)[:200]
                })
            if attempt < _URL_RETRY_ATTEMPTS:
                actlog.log(EVT_TOOL_RETRY, {
                    "url": url, "attempt": attempt, "reason": type(exc).__name__, "backoff_s": _URL_RETRY_BACKOFF_BASE_SECS * (2 ** (attempt - 1))
                })
                delay = _URL_RETRY_BACKOFF_BASE_SECS * (2 ** (attempt - 1))
                time.sleep(delay)
                continue
            raise
        except httpx.HTTPStatusError as exc:
            last_exc = exc
            raise

    if last_exc:
        raise last_exc
    raise RuntimeError("visit_url: unexpected fetch termination")


# ── Per-tool rate limiters ────────────────────────────────────────────────────

class _RateLimiter:
    """Enforce a minimum interval between successive calls (blocks the calling thread)."""

    def __init__(self, min_interval_secs: float) -> None:
        self._lock = threading.Lock()
        self._last_call: float = 0.0
        self._min_interval = min_interval_secs

    def wait(self) -> None:
        with self._lock:
            now = time.monotonic()
            gap = self._min_interval - (now - self._last_call)
            if gap > 0:
                time.sleep(gap)
            self._last_call = time.monotonic()


_rate_web_search = _RateLimiter(2.0)   # ~1 call/2s for DuckDuckGo
_rate_news_search = _RateLimiter(2.0)  # ~1 call/2s for DuckDuckGo news
_rate_arxiv = _RateLimiter(3.0)        # ~1 call/3s per arXiv guidelines


def set_runtime_context(user_id: str, workspace_id: str, session_id: str, run_id: str = "") -> None:
    _tls.user_id = user_id
    _tls.workspace_id = workspace_id
    _tls.session_id = session_id
    _tls.run_id = run_id
    actlog.set_log_context(user_id, session_id, run_id)


def _tool_settings() -> dict:
    user_id = getattr(_tls, "user_id", "")
    if not user_id:
        return get_user_tool_settings("") if False else {
            "web_search": {"default_max_results": 4, "max_results_limit": 10},
            "news_search": {"default_max_results": 8, "max_results_limit": 15},
            "arxiv": {"default_max_results": 5, "max_results_limit": 15},
            "youtube_transcript": {"default_max_videos": 3, "max_videos_limit": 10, "skip_auto_generated": False},
        }
    return get_user_tool_settings(user_id)


def _session_output_dir() -> Path:
    user_id = getattr(_tls, "user_id", "user")
    workspace_id = getattr(_tls, "workspace_id", "workspace")
    session_id = getattr(_tls, "session_id", "session")
    root = Path(__file__).resolve().parents[2] / "instance" / "files"
    folder = root / user_id[:8] / workspace_id[:8] / session_id[:8]
    folder.mkdir(parents=True, exist_ok=True)
    return folder


def _slide_structuring_prompt(min_slides: int, max_slides: int) -> str:
    return (
        "You structure research text into slide-ready content.\n"
        f"Target at least {min_slides} and up to {max_slides} slides, but use more if needed for completeness.\n"
        "Return plain text with this shape only:\n"
        "# Slide Title\n"
        "- bullet\n"
        "- bullet\n\n"
        "Rules:\n"
        "- Keep every bullet concise but information-dense.\n"
        "- Preserve key facts and numbers.\n"
        "- No markdown tables, no code fences, no JSON."
    )


def _normalize_llm_content(content) -> str:
    if isinstance(content, list):
        return " ".join(
            part.get("text", "") if isinstance(part, dict) else str(part)
            for part in content
        ).strip()
    return str(content or "").strip()


def _reformat_for_slides(raw_text: str) -> str:
    user_id = getattr(_tls, "user_id", "")
    workspace_id = getattr(_tls, "workspace_id", "")
    if not user_id or not workspace_id:
        return raw_text

    slide_cfg = (_tool_settings().get("slides") or {})
    min_slides = int(slide_cfg.get("min_slides", 5) or 5)
    max_slides = int(slide_cfg.get("max_slides", 10) or 10)

    llm_chain = build_summary_llms(user_id, workspace_id)
    if not llm_chain:
        return raw_text

    override = str(slide_cfg.get("model_override", "") or "").strip()
    if override:
        preferred = [entry for entry in llm_chain if entry[1] == override]
        fallback = [entry for entry in llm_chain if entry[1] != override]
        llm_chain = preferred + fallback

    messages = [
        SystemMessage(content=_slide_structuring_prompt(min_slides=min_slides, max_slides=max_slides)),
        HumanMessage(content=raw_text),
    ]

    for llm, model_name in llm_chain:
        try:
            response = llm.invoke(messages)
            structured = _normalize_llm_content(getattr(response, "content", ""))
            if structured:
                actlog.log("slide_structured", {"model": model_name, "chars": len(structured)})
                return structured
        except Exception as exc:
            actlog.log("slide_structuring_retry", {"model": model_name, "error": str(exc)[:300]})
            continue

    actlog.log("slide_structuring_fallback_raw", {"reason": "all_models_failed"})
    return raw_text


_TEMPLATE_SECTIONS: dict[str, str] = {
    "research_report": (
        "1. Title (use the provided title or derive one)\n"
        "2. Executive Summary\n"
        "3. Key Findings\n"
        "4. Detailed Analysis\n"
        "5. Sources / References (if any citations/links are present)\n"
        "6. Missing / Not specified (if information is lacking)"
    ),
    "procurement_report": (
        "1. Title\n"
        "2. Objective\n"
        "3. Recommended Approach\n"
        "4. Bill of Materials (use a Markdown table)\n"
        "5. Vendor / Seller Options (use a Markdown table if available)\n"
        "6. Cost Notes\n"
        "7. Safety / Compatibility Notes\n"
        "8. Sources / References (if present)\n"
        "9. Missing / Not specified (if needed)"
    ),
    "comparison_report": (
        "1. Title\n"
        "2. Decision Summary\n"
        "3. Comparison Table (use a Markdown table)\n"
        "4. Option-by-option Analysis\n"
        "5. Tradeoffs\n"
        "6. Recommendation\n"
        "7. Sources / References (if present)\n"
        "8. Missing / Not specified (if needed)"
    ),
    "financial_report": (
        "1. Title\n"
        "2. Executive Summary\n"
        "3. Key Metrics\n"
        "4. Analysis\n"
        "5. Risks / Assumptions\n"
        "6. Sources / References (if present)\n"
        "7. Missing / Not specified (if needed)"
    ),
}

_VALID_TEMPLATES = set(_TEMPLATE_SECTIONS.keys())
_VALID_STYLES = {"clean", "dense", "executive"}


def _pdf_structuring_prompt(template: str, title: str = "") -> str:
    sections = _TEMPLATE_SECTIONS.get(template, _TEMPLATE_SECTIONS["research_report"])
    prompt = (
        "You are a professional report writer. Transform the input below into a polished "
        "Markdown report.\n\n"
        "Requirements:\n"
        "- Use clear headings (## or ###).\n"
        "- Use Markdown tables where useful for structured data.\n"
        "- Preserve any links and citations already present in the input.\n"
        "- Do NOT invent sources or citations.\n"
        "- If key information is missing, include a section called 'Missing / Not specified'\n"
        "  indicating what is absent rather than fabricating content.\n"
        "- Avoid code fences unless the source data specifically requires code.\n"
        "- Output Markdown only. Do not include surrounding explanation.\n\n"
    )
    if title:
        prompt += f"Report title: {title}\n\n"
    prompt += "Required report structure:\n"
    prompt += sections
    return prompt


def _format_for_pdf(data: str, template: str, title: str) -> str:
    user_id = getattr(_tls, "user_id", "")
    workspace_id = getattr(_tls, "workspace_id", "")
    if not user_id or not workspace_id:
        return data

    llm_chain = build_summary_llms(user_id, workspace_id)
    if not llm_chain:
        return data

    messages = [
        SystemMessage(content=_pdf_structuring_prompt(template, title)),
        HumanMessage(content=data),
    ]

    try:
        response, model_name = invoke_with_fallbacks(llm_chain, messages)
        structured = _normalize_llm_content(getattr(response, "content", ""))
        if structured:
            actlog.log("pdf_structured", {
                "model": model_name, "chars": len(structured), "template": template
            })
            return structured
    except Exception as exc:
        actlog.log("pdf_structuring_fallback_raw", {"error": str(exc)[:300]})

    return data


def _safe_artifact_filename(filename: str, extension: str) -> str:
    name = re.sub(r'[^\w\-_\. ]', '', str(filename or "report"))
    name = name.strip()
    if not name:
        name = "report"
    if not name.endswith(extension):
        name = f"{name}{extension}"
    return name


def _markdown_to_html(markdown_text: str) -> str:
    return _md_lib.markdown(
        markdown_text,
        extensions=["extra", "tables", "fenced_code", "sane_lists"],
    )


def _build_report_html(body_html: str, title: str, template: str, style: str) -> str:
    ts = datetime.utcnow().strftime("%Y-%m-%d %H:%M UTC")
    style_class = style if style in _VALID_STYLES else "clean"
    disclaimer = ""
    if template == "financial_report":
        disclaimer = (
            '<div class="disclaimer">'
            "This report is generated for informational purposes only "
            "and is not financial advice."
            "</div>"
        )

    escaped_title = title.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")

    css = """
@page {
  size: A4;
  margin: 2cm 2.5cm;
}

body {
  font-family: 'Helvetica Neue', Helvetica, Arial, sans-serif;
  font-size: 11pt;
  line-height: 1.5;
  color: #1a1a1a;
}

h1 { font-size: 20pt; margin-top: 0; margin-bottom: 0.3cm; color: #111; }
h2 { font-size: 14pt; margin-top: 0.6cm; margin-bottom: 0.3cm;
      color: #222; border-bottom: 1px solid #ccc; padding-bottom: 0.1cm; }
h3 { font-size: 12pt; margin-top: 0.4cm; margin-bottom: 0.2cm; color: #333; }
h4 { font-size: 11pt; margin-top: 0.3cm; margin-bottom: 0.2cm; color: #444; }

p { margin: 0.2cm 0; }
ul, ol { margin: 0.2cm 0; padding-left: 1.2cm; }
li { margin: 0.1cm 0; }

table {
  width: 100%;
  border-collapse: collapse;
  margin: 0.4cm 0;
  font-size: 10pt;
  page-break-inside: auto;
}

th, td {
  border: 1px solid #999;
  padding: 6px 8px;
  text-align: left;
  vertical-align: top;
}

th {
  background-color: #eef;
  font-weight: bold;
}

tr:nth-child(even) td {
  background-color: #f8f8f8;
}

thead { display: table-header-group; }

.metadata {
  font-size: 9pt;
  color: #666;
  margin-bottom: 0.5cm;
  padding-bottom: 0.3cm;
  border-bottom: 1px solid #ddd;
}

.disclaimer {
  margin-top: 1cm;
  padding: 0.3cm 0.5cm;
  border: 1px solid #cc0000;
  background-color: #fff5f5;
  font-size: 9pt;
  color: #800;
  text-align: center;
}

a { color: #2563eb; text-decoration: underline; }
code { font-family: 'Courier New', monospace; font-size: 10pt;
       background-color: #f4f4f4; padding: 1px 4px; border-radius: 2px; }
pre { background-color: #f4f4f4; padding: 0.3cm; border: 1px solid #ddd;
      border-radius: 3px; overflow-x: auto; font-size: 9pt; }
blockquote { border-left: 3px solid #ccc; margin: 0.3cm 0; padding: 0.1cm 0.5cm;
             color: #555; }

/* dense */
.dense body { font-size: 9.5pt; line-height: 1.3; }
.dense h1 { font-size: 16pt; }
.dense h2 { font-size: 12pt; margin-top: 0.4cm; }
.dense h3 { font-size: 10.5pt; margin-top: 0.3cm; }
.dense p { margin: 0.1cm 0; }
.dense table { font-size: 8.5pt; }
.dense th, .dense td { padding: 4px 6px; }
.dense ul, .dense ol { margin: 0.1cm 0; }

/* executive */
.executive body { font-size: 12pt; line-height: 1.6; }
.executive h1 { font-size: 24pt; margin-bottom: 0.5cm; }
.executive h2 { font-size: 16pt; border-bottom-width: 2px; }
.executive h3 { font-size: 13pt; }
.executive p { margin: 0.3cm 0; }
.executive .metadata { margin-bottom: 0.8cm; }
"""

    return f"""<!DOCTYPE html>
<html lang="en">
<head>
<meta charset="utf-8">
<title>{escaped_title}</title>
<style>{css}</style>
</head>
<body class="{style_class}">
<h1>{escaped_title or "Report"}</h1>
<div class="metadata">Template: {template} &nbsp;|&nbsp; Style: {style} &nbsp;|&nbsp; Generated: {ts}</div>
{body_html}
{disclaimer}
</body>
</html>"""


@tool("web_search")
def web_search(query: str, max_results: int = 4) -> str:
    """Search the web using DuckDuckGo and return title/link/snippet results."""
    cfg = _tool_settings().get("web_search", {})
    default_max = int(cfg.get("default_max_results", 4))
    limit = int(cfg.get("max_results_limit", 10))
    effective = min(max_results if max_results != 4 else default_max, limit)
    _rate_web_search.wait()
    try:
        with DDGS() as ddgs:
            rows = list(ddgs.text(query, max_results=effective))
    except Exception as exc:
        return f"Error searching web: {exc}"

    if not rows:
        return "No search results found."
    lines = [f"Query: {query}", f"Results: {len(rows)}"]
    for index, row in enumerate(rows, 1):
        lines.append(f"{index}. {row.get('title', '(no title)')}")
        link = row.get("href") or row.get("url") or row.get("link")
        snippet = row.get("body") or row.get("snippet")
        if link:
            lines.append(f"   link: {link}")
        if snippet:
            lines.append(f"   snippet: {snippet}")
    return "\n".join(lines)


@tool("wiki_search")
def wiki_search(query: str) -> str:
    """Search Wikipedia and return a concise result snippet."""
    return _wiki.run(query)


@tool("visit_url")
def visit_url(url: str, format: str = "text") -> str:
    """Fetch a URL and return cleaned text, HTML, or JSON output."""
    if not re.match(r"^https?://", url or ""):
        return "Error: URL must start with http:// or https://"

    try:
        response = _fetch_url_with_retries(url)
    except Exception as exc:
        return f"Error fetching URL: {exc}"

    if format == "html":
        return response.text[:_URL_MAX_CHARS]

    soup = BeautifulSoup(response.text, "html.parser")
    for tag in soup(["script", "style", "noscript", "footer", "nav", "aside", "iframe", "form"]):
        tag.decompose()

    title = soup.title.string.strip() if soup.title and soup.title.string else "Untitled"
    body = re.sub(r"\n{3,}", "\n\n", soup.get_text("\n", strip=True))[:_URL_MAX_CHARS]

    if format == "json":
        return json.dumps({"url": url, "title": title, "text": body}, ensure_ascii=False)

    return f"--- Page: {title} ---\nURL: {url}\n\n{body}"


@tool("news_search")
def news_search(query: str, max_results: int = 8) -> str:
    """Search news headlines using DuckDuckGo news results."""
    cfg = _tool_settings().get("news_search", {})
    default_max = int(cfg.get("default_max_results", 8))
    limit = int(cfg.get("max_results_limit", 15))
    effective = min(max_results if max_results != 8 else default_max, limit)
    _rate_news_search.wait()
    try:
        with DDGS() as ddgs:
            rows = list(ddgs.news(query, max_results=effective))
    except Exception as exc:
        return f"Error searching news: {exc}"

    if not rows:
        return "No news results found."

    lines = [f"News query: {query}"]
    for index, row in enumerate(rows, 1):
        lines.append(f"{index}. {row.get('title', '(no title)')}")
        if row.get("source"):
            lines.append(f"   source: {row['source']}")
        if row.get("date"):
            lines.append(f"   date: {row['date']}")
        if row.get("url"):
            lines.append(f"   url: {row['url']}")
    return "\n".join(lines)


def _extract_video_id(url_or_id: str) -> str | None:
    match = re.search(r"(?:v=|youtu\.be/|embed/|shorts/)([A-Za-z0-9_-]{11})", url_or_id)
    return match.group(1) if match else None


@tool("youtube_transcript")
def youtube_transcript(query_or_url: str, max_videos: int = 3) -> str:
    """Fetch transcripts from direct YouTube URLs or YouTube search results."""
    cfg = _tool_settings().get("youtube_transcript", {})
    default_max = int(cfg.get("default_max_videos", 3))
    limit = int(cfg.get("max_videos_limit", 10))
    skip_auto = bool(cfg.get("skip_auto_generated", False))
    effective = min(max_videos if max_videos != 3 else default_max, limit)

    candidates: list[dict] = []
    direct_id = _extract_video_id(query_or_url)
    if direct_id:
        candidates.append({"id": direct_id, "title": query_or_url, "url": query_or_url})
    else:
        try:
            with DDGS() as ddgs:
                videos = list(ddgs.videos(f"{query_or_url} site:youtube.com", max_results=effective * 3))
        except Exception as exc:
            return f"Error searching videos: {exc}"

        for item in videos:
            url = item.get("content") or item.get("url") or ""
            video_id = _extract_video_id(url)
            if video_id:
                candidates.append({"id": video_id, "title": item.get("title") or url, "url": url})
            if len(candidates) >= effective:
                break

    if not candidates:
        return "No YouTube videos found."

    chunks: list[str] = []
    for candidate in candidates:
        vid = candidate["id"]
        title = candidate["title"]
        url = candidate["url"]
        try:
            transcript_list = YouTubeTranscriptApi.list_transcripts(vid)
            used_auto = False
            try:
                transcript = transcript_list.find_manually_created_transcript(["en"])
            except NoTranscriptFound:
                transcript = transcript_list.find_generated_transcript(["en"])
                used_auto = True

            if used_auto and skip_auto:
                chunks.append(f"[Video: {title}]\nURL: {url}\nSkipped auto-generated transcript.\n")
                continue

            rows = transcript.fetch()
            text = " ".join((row["text"] if isinstance(row, dict) else row.text) for row in rows)
            chunks.append(f"[Video: {title}]\nURL: {url}\n\n{text}\n")
        except TranscriptsDisabled:
            chunks.append(f"[Video: {title}]\nURL: {url}\nTranscripts are disabled.\n")
        except NoTranscriptFound:
            chunks.append(f"[Video: {title}]\nURL: {url}\nNo English transcript found.\n")
        except Exception as exc:
            chunks.append(f"[Video: {title}]\nURL: {url}\nError: {exc}\n")

    return "\n---\n".join(chunks)


@tool("arxiv_search")
def arxiv_search(query: str, max_results: int = 5) -> str:
    """Search arXiv and return structured paper summaries."""
    cfg = _tool_settings().get("arxiv", {})
    default_max = int(cfg.get("default_max_results", 5))
    limit = int(cfg.get("max_results_limit", 15))
    effective = min(max_results if max_results != 5 else default_max, limit)

    _rate_arxiv.wait()
    try:
        client = arxiv.Client()
        search = arxiv.Search(query=query, max_results=effective, sort_by=arxiv.SortCriterion.Relevance)
        rows = list(client.results(search))
    except Exception as exc:
        return f"Error searching arXiv: {exc}"

    if not rows:
        return "No arXiv papers found."

    lines = [f"arXiv query: {query}"]
    for index, paper in enumerate(rows, 1):
        authors = ", ".join(author.name for author in paper.authors[:4])
        if len(paper.authors) > 4:
            authors += " et al."
        lines.extend(
            [
                f"{index}. {paper.title}",
                f"   authors: {authors}",
                f"   published: {paper.published.strftime('%Y-%m-%d') if paper.published else 'N/A'}",
                f"   link: {paper.entry_id}",
                f"   abstract: {paper.summary[:500]}",
            ]
        )
    return "\n".join(lines)


@tool("save_text")
def save_text(data: str, filename: str = "research_output.txt") -> str:
    """Save text to a session-scoped output file."""
    folder = _session_output_dir()
    target = Path(filename)
    if not target.is_absolute():
        target = folder / target.name

    timestamp = datetime.utcnow().isoformat()
    with open(target, "a", encoding="utf-8") as handle:
        handle.write(f"--- {timestamp} ---\n{data}\n\n")
    return f"Saved text to {target}"


@tool("create_slides")
def create_slides(data: str, format: str = "ppt", filename: str = "slides") -> str:
    """Create simple slide artifacts from text as PPT or image scaffold output."""
    folder = _session_output_dir()
    fmt = (format or "ppt").strip().lower()
    structured_data = _reformat_for_slides(data)

    if fmt == "images":
        target_dir = folder / filename
        target_dir.mkdir(parents=True, exist_ok=True)
        image_stub = target_dir / "slide_01.txt"
        with open(image_stub, "w", encoding="utf-8") as handle:
            handle.write("Image slide scaffold\n\n")
            handle.write(structured_data)
        return f"Slide image scaffold saved to {target_dir}"

    ppt_path = folder / (filename if filename.endswith(".pptx") else f"{filename}.pptx")
    try:
        from pptx import Presentation

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Research Slides"
        body = slide.placeholders[1].text_frame
        lines = [line.strip() for line in structured_data.splitlines() if line.strip()][:8]
        if not lines:
            lines = ["No content provided."]
        body.text = lines[0]
        for line in lines[1:]:
            paragraph = body.add_paragraph()
            paragraph.text = line
        prs.save(str(ppt_path))
        return f"Slide deck saved to {ppt_path}"
    except Exception:
        fallback = folder / f"{filename}.txt"
        with open(fallback, "w", encoding="utf-8") as handle:
            handle.write(structured_data)
        return f"PPT generation unavailable; slide text saved to {fallback}"


@tool("create_pdf")
def create_pdf(
    data: str,
    template: str = "research_report",
    filename: str = "report",
    title: str = "",
    style: str = "clean",
) -> str:
    """Generate a polished PDF report from research findings.

    Use this tool when the user specifically requests a PDF document, a formal
    written report, or a downloadable summary of the research session.

    Args:
        data: Research findings / conversation summary to convert into the PDF.
        template: Report layout. One of: research_report, procurement_report,
                  comparison_report, financial_report. Default: research_report.
        filename: Base name for the output file (without extension).
                  Default: "report".
        title: Optional report title shown in the PDF heading.
        style: Visual appearance. One of: clean, dense, executive.
               Default: clean.

    Notes:
        - This tool does NOT perform web searches. Research first, then call
          this tool with the gathered data.
        - Allowed templates: research_report, procurement_report,
          comparison_report, financial_report.
    """
    valid_templates = {"research_report", "procurement_report", "comparison_report", "financial_report"}
    valid_styles = {"clean", "dense", "executive"}

    effective_template = template if template in valid_templates else "research_report"
    effective_style = style if style in valid_styles else "clean"

    sanitized = _safe_artifact_filename(filename, ".pdf")

    folder = _session_output_dir()
    pdf_path = folder / sanitized
    html_path = folder / _safe_artifact_filename(filename, ".html")

    try:
        report_md = _format_for_pdf(data, effective_template, title)
    except Exception as exc:
        report_md = data

    try:
        body_html = _markdown_to_html(report_md)
    except Exception as exc:
        return f"PDF generation failed: error rendering Markdown: {exc}"

    full_html = _build_report_html(body_html, title, effective_template, effective_style)

    try:
        from weasyprint import HTML as _WeasyHTML

        _WeasyHTML(string=full_html).write_pdf(str(pdf_path))
        return f"PDF report saved to {pdf_path}"
    except Exception as exc:
        try:
            with open(html_path, "w", encoding="utf-8") as f:
                f.write(full_html)
            return (
                f"PDF generation failed; HTML fallback saved to {html_path}. "
                f"Error: {exc}"
            )
        except Exception as inner:
            return (
                f"PDF generation failed and HTML fallback also failed. "
                f"Error: {exc}; Fallback error: {inner}"
            )


ALL_TOOLS = [
    web_search,
    wiki_search,
    visit_url,
    news_search,
    youtube_transcript,
    arxiv_search,
    save_text,
    create_slides,
    create_pdf,
]

TOOL_MAP = {tool_item.name: tool_item for tool_item in ALL_TOOLS}
